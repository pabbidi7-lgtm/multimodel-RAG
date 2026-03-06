"""
Multimodal RAG API v3
=====================
Endpoints : /create  /ingest  /ask
Formats   : PDF · DOCX · PPTX · HTML · TXT
Pipeline  : Parse → Chunk → OCR → Batch Embed → Milvus Lite → Rerank → LLM → Confidence
"""

import os, re, json, uuid, time, base64, hashlib, datetime, socket, logging
import numpy as np
import requests
import fitz
import cv2

from io import BytesIO
from contextlib import asynccontextmanager
from dataclasses import dataclass, field
from typing import Optional

from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from requests.adapters import HTTPAdapter
from urllib3.util.retry import Retry
from fastapi import FastAPI, UploadFile, File, HTTPException, Query
from pymilvus import MilvusClient, DataType
from dotenv import load_dotenv

# ── PaddleOCR (optional) ──────────────────────────────────────────────────────
try:
    from paddleocr import PaddleOCR
    _ocr = PaddleOCR(use_angle_cls=True, lang="en", show_log=False)
    OCR_AVAILABLE = True
except Exception:
    _ocr = None
    OCR_AVAILABLE = False

# =============================================================================
# LOGGING
# =============================================================================

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s  [%(levelname)-8s]  %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)
log = logging.getLogger("rag")

# =============================================================================
# CONFIG
# =============================================================================

load_dotenv()

NVIDIA_KEY = os.getenv("NVIDIA_API")
if not NVIDIA_KEY:
    raise RuntimeError("NVIDIA_API not found in .env")

BASE_URL     = "https://integrate.api.nvidia.com/v1"
EMBED_MODEL  = "nvidia/llama-3.2-nemoretriever-1b-vlm-embed-v1"
RERANK_MODEL = "nvidia/llama-3.2-nemoretriever-1b-vlm-embed-v1"
PRIMARY_LLM  = "meta/llama3-70b-instruct"
FALLBACK_LLM = "meta/llama3-8b-instruct"

COLLECTION   = "rag_store"
DB_PATH      = "./rag.db"
CHUNK_TOKENS = 400
OVERLAP      = 0.15
MAX_CONTENT  = 4800
TOKEN_BUDGET = 3000
DEDUP_THRESH = 0.95
EMBED_BATCH  = 32        # safe batch size — avoids NVIDIA modality-length errors

# =============================================================================
# HTTP SESSION  (IPv4 forced — fixes silent 10-min hang on broken-IPv6 Linux)
# =============================================================================

class _IPv4Adapter(HTTPAdapter):
    def send(self, *args, **kwargs):
        _orig = socket.getaddrinfo
        def _v4(host, port, family=0, *a, **kw):
            return _orig(host, port, socket.AF_INET, *a, **kw)
        socket.getaddrinfo = _v4
        try:
            return super().send(*args, **kwargs)
        finally:
            socket.getaddrinfo = _orig

def _make_session():
    s = requests.Session()
    retry = Retry(total=2, backoff_factor=1.0,
                  status_forcelist=[429, 500, 502, 503, 504],
                  allowed_methods=["POST", "GET"], connect=0, read=2)
    s.mount("https://", _IPv4Adapter(max_retries=retry))
    s.mount("http://",  _IPv4Adapter(max_retries=retry))
    return s

session = _make_session()

# =============================================================================
# DATA MODEL
# =============================================================================

@dataclass
class Chunk:
    text:      str
    modality:  str        # text | table | image_ocr
    doc_id:    str
    file_name: str
    file_type: str
    page:      int = 0
    section:   str = ""
    parent_id: str = ""

# =============================================================================
# UTILITIES
# =============================================================================

def normalize(v):
    a = np.array(v, dtype=np.float32)
    n = np.linalg.norm(a)
    return (a / n).tolist() if n > 0 else a.tolist()

def token_count(t: str) -> int:
    return len(t.split())

def cosine(a, b) -> float:
    a, b = np.array(a), np.array(b)
    d = np.linalg.norm(a) * np.linalg.norm(b)
    return float(np.dot(a, b) / d) if d > 0 else 0.0

def now() -> str:
    return datetime.datetime.utcnow().isoformat()

# =============================================================================
# CHUNKING
# =============================================================================

def chunk_text(text: str) -> list:
    sentences = [s.strip() for s in text.replace("\n", " ").split(". ") if s.strip()]
    if not sentences:
        return [text[:MAX_CONTENT]] if text.strip() else []
    chunks, buf = [], []
    for sent in sentences:
        buf.append(sent)
        if token_count(". ".join(buf)) >= CHUNK_TOKENS:
            chunks.append(". ".join(buf) + ".")
            keep = max(1, int(len(buf) * OVERLAP))
            buf  = buf[-keep:]
    if buf:
        chunks.append(". ".join(buf) + ".")
    return chunks or [text[:MAX_CONTENT]]

# =============================================================================
# OCR
# =============================================================================

def run_ocr(img_bytes: bytes) -> str:
    if not OCR_AVAILABLE or not img_bytes:
        return ""
    try:
        arr    = np.frombuffer(img_bytes, dtype=np.uint8)
        img    = cv2.imdecode(arr, cv2.IMREAD_COLOR)
        result = _ocr.ocr(img, cls=True)
        if not result:
            return ""
        return " ".join(line[1][0] for block in result if block for line in block)
    except Exception as e:
        log.warning(f"OCR error: {e}")
        return ""

# =============================================================================
# PARSERS
# =============================================================================

def _table_text(rows):
    if not rows:
        return ""
    header = " | ".join(rows[0])
    body   = "\n".join(" | ".join(r) for r in rows[1:])
    return f"{header}\n{'-'*len(header)}\n{body}"

def _img_chunk(img_bytes, caption, doc_id, fname, ftype, page, section):
    """
    Create a Chunk for an image.
    Priority: OCR text > figure caption > placeholder.
    Always returns a chunk (never None) so figure captions are always indexed.
    """
    pid = str(uuid.uuid4())[:12]

    # Try OCR first (only if image bytes available and OCR is on)
    ocr_text = ""
    if img_bytes and len(img_bytes) >= 5000:
        ocr_text = run_ocr(img_bytes).strip()

    if ocr_text:
        # Best case: OCR extracted real text from the image
        text = f"{caption} | OCR: {ocr_text}" if caption else ocr_text
        log.debug(f"Image OCR ({ftype} p{page}): {len(ocr_text)} chars")
    elif caption:
        # No OCR but we have the figure caption — use it as the chunk text.
        # This makes figures queryable even when OCR is off or image is vector.
        text = caption
        log.debug(f"Image no-OCR ({ftype} p{page}): using caption '{caption[:60]}'")
    else:
        # Last resort placeholder
        text = f"[Figure on page {page}]"

    return Chunk(text, "image_ocr", doc_id, fname, ftype, page, section, pid)


def _extract_fig_captions(page) -> dict:
    """
    Extract all figure captions on a page.
    Returns dict of {fig_num_str: full_caption} e.g. {"1": "Fig. 1. The safety ..."}
    Also returns a list of all captions for pages where fig number is unclear.
    """
    captions = {}
    caption_list = []
    for block in page.get_text("blocks"):
        if block[6] != 0:
            continue
        txt = block[4].strip().replace("\n", " ")
        if re.match(r"^[Ff]ig\.?\s*\d+", txt):
            caption_list.append(txt)
            m = re.match(r"^[Ff]ig\.?\s*(\d+)", txt)
            if m:
                captions[m.group(1)] = txt
    return captions, caption_list


def parse_pdf(raw, fname, doc_id):
    log.info(f"Parsing PDF: {fname}")
    chunks, pdf = [], fitz.open(stream=raw, filetype="pdf")
    heading = ""
    seen_xrefs = set()     # track xrefs so same image embedded on 2 pages isn't embedded twice

    for pnum, page in enumerate(pdf, 1):

        # ── Step 1: Extract figure captions first ──────────────────────────
        fig_captions_dict, fig_captions_list = _extract_fig_captions(page)
        if fig_captions_list:
            log.debug(f"Page {pnum} captions: {fig_captions_list}")

        # ── Step 2: Parse text blocks ───────────────────────────────────────
        parts = []
        for block in page.get_text("blocks"):
            if block[6] != 0:
                continue
            txt = block[4].strip().replace("\n", " ")
            if not txt:
                continue

            # Section heading detection: short, ALL-CAPS
            if len(txt) < 80 and txt.isupper():
                if parts:
                    for c in chunk_text(" ".join(parts)):
                        chunks.append(Chunk(c, "text", doc_id, fname, "pdf", pnum, heading))
                    parts = []
                heading = txt
                continue

            # Skip figure caption lines from body text — they'll be stored with images
            if re.match(r"^[Ff]ig\.?\s*\d+", txt):
                continue

            parts.append(txt)

        if parts:
            for c in chunk_text(" ".join(parts)):
                chunks.append(Chunk(c, "text", doc_id, fname, "pdf", pnum, heading))

        # ── Step 3: Raster images via fitz ────────────────────────────────
        page_img_count = 0
        for img_info in page.get_images(full=True):
            xref = img_info[0]
            if xref in seen_xrefs:
                log.debug(f"Page {pnum}: skipping duplicate xref={xref}")
                continue
            seen_xrefs.add(xref)

            try:
                img_bytes = pdf.extract_image(xref)["image"]
                if len(img_bytes) < 5000:
                    log.debug(f"Page {pnum}: xref={xref} skipped (<5KB icon)")
                    continue
                # Pick the best caption for this image (use first available on the page)
                caption = fig_captions_list[page_img_count] if page_img_count < len(fig_captions_list) else (fig_captions_list[0] if fig_captions_list else "")
                ch = _img_chunk(img_bytes, caption, doc_id, fname, "pdf", pnum, heading)
                chunks.append(ch)
                page_img_count += 1
                log.debug(f"Page {pnum}: raster image xref={xref} → chunk with caption='{caption[:50]}'")
            except Exception as e:
                log.warning(f"PDF image p{pnum} xref={xref}: {e}")

        # ── Step 4: Vector figures (Form XObjects with no raster bytes) ────
        # MATLAB/LaTeX often export figures as pure vector PDFs.
        # fitz.get_images() misses these. We detect them via fig captions
        # that weren't matched to any raster image above.
        if fig_captions_list and page_img_count == 0:
            # This page has figure captions but no raster images → vector figure
            for caption in fig_captions_list:
                ch = _img_chunk(None, caption, doc_id, fname, "pdf", pnum, heading)
                chunks.append(ch)
                log.debug(f"Page {pnum}: vector figure chunk from caption '{caption[:60]}'")

    log.info(f"PDF '{fname}': {len(chunks)} chunks, {len(pdf)} pages")
    log.info(f"  text={sum(1 for c in chunks if c.modality=='text')} "
             f"image_ocr={sum(1 for c in chunks if c.modality=='image_ocr')}")
    return chunks


def parse_docx(raw, fname, doc_id):
    log.info(f"Parsing DOCX: {fname}")
    chunks, doc, heading = [], Document(BytesIO(raw)), ""

    for para in doc.paragraphs:
        t = para.text.strip()
        if not t:
            continue
        if para.style.name.startswith("Heading"):
            heading = t
        else:
            for c in chunk_text(t):
                chunks.append(Chunk(c, "text", doc_id, fname, "docx", 0, heading))

    for table in doc.tables:
        rows  = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        ttext = _table_text(rows)
        if ttext:
            chunks.append(Chunk(ttext, "table", doc_id, fname, "docx", 0, heading))

    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                ch = _img_chunk(rel.target_part.blob, "", doc_id, fname, "docx", 0, heading)
                chunks.append(ch)
            except Exception:
                pass

    log.info(f"DOCX '{fname}': {len(chunks)} chunks")
    return chunks


def parse_pptx(raw, fname, doc_id):
    log.info(f"Parsing PPTX: {fname}")
    chunks, prs = [], Presentation(BytesIO(raw))

    for i, slide in enumerate(prs.slides, 1):
        sec = f"Slide {i}"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                for c in chunk_text(shape.text.strip()):
                    chunks.append(Chunk(c, "text", doc_id, fname, "pptx", i, sec))
            if shape.shape_type == 13:
                try:
                    ch = _img_chunk(shape.image.blob, "", doc_id, fname, "pptx", i, sec)
                    chunks.append(ch)
                except Exception:
                    pass
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                chunks.append(Chunk(notes, "text", doc_id, fname, "pptx", i, f"Slide {i} Notes"))

    log.info(f"PPTX '{fname}': {len(chunks)} chunks, {len(prs.slides)} slides")
    return chunks


def parse_html(raw, fname, doc_id):
    log.info(f"Parsing HTML: {fname}")
    chunks = []
    soup   = BeautifulSoup(raw.decode(errors="ignore"), "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    heading = ""
    for tag in soup.find_all(["h1","h2","h3","p","table"]):
        if tag.name in ["h1","h2","h3"]:
            heading = tag.get_text(strip=True)
        elif tag.name == "p":
            t = tag.get_text(strip=True)
            if t:
                for c in chunk_text(t):
                    chunks.append(Chunk(c, "text", doc_id, fname, "html", 0, heading))
        elif tag.name == "table":
            rows  = [[td.get_text(strip=True) for td in tr.find_all(["td","th"])]
                     for tr in tag.find_all("tr")]
            ttext = _table_text([r for r in rows if r])
            if ttext:
                chunks.append(Chunk(ttext, "table", doc_id, fname, "html", 0, heading))
    log.info(f"HTML '{fname}': {len(chunks)} chunks")
    return chunks


def parse_file(raw, fname, doc_id):
    ext = fname.rsplit(".", 1)[-1].lower()
    fn  = {"pdf": parse_pdf, "docx": parse_docx, "pptx": parse_pptx,
           "html": parse_html, "htm": parse_html}.get(ext)
    if fn:
        return fn(raw, fname, doc_id)
    # Plain text fallback
    log.info(f"Parsing TXT: {fname}")
    return [Chunk(c, "text", doc_id, fname, "txt", 0, "")
            for c in chunk_text(raw.decode(errors="ignore"))]

# =============================================================================
# EMBEDDING
# KEY FIX: "modality" list length MUST equal "input" list length
# =============================================================================

def _post(endpoint, payload, timeout=(5, 90)):
    log.debug(f"POST /{endpoint}")
    try:
        r = session.post(
            f"{BASE_URL}/{endpoint}",
            headers={"Authorization": f"Bearer {NVIDIA_KEY}"},
            json=payload, timeout=timeout,
        )
    except requests.exceptions.ConnectionError as e:
        log.error(f"Connection error → {endpoint}: {e}")
        raise HTTPException(503, f"Cannot reach NVIDIA API: {e}")
    except requests.exceptions.Timeout:
        log.error(f"Timeout → {endpoint}")
        raise HTTPException(504, "NVIDIA API timed out")

    if r.status_code != 200:
        log.error(f"NVIDIA {endpoint} → {r.status_code}: {r.text[:400]}")
        raise HTTPException(502, f"NVIDIA API {endpoint} error {r.status_code}: {r.text[:300]}")

    return r.json()


def embed_texts(texts: list, mode="passage") -> list:
    """
    Batch embed. CRITICAL: modality list must match input list length exactly.
    e.g. 32 texts → "modality": ["text", "text", ..., "text"]  (32 times)
    """
    if not texts:
        return []
    all_vecs    = []
    n_batches   = (len(texts) + EMBED_BATCH - 1) // EMBED_BATCH

    for i in range(0, len(texts), EMBED_BATCH):
        batch     = texts[i:i + EMBED_BATCH]
        batch_num = i // EMBED_BATCH + 1
        log.info(f"Embedding batch {batch_num}/{n_batches} — {len(batch)} texts")

        data = _post("embeddings", {
            "model":           EMBED_MODEL,
            "input":           batch,
            "input_type":      mode,
            "encoding_format": "float",
            "modality":        ["text"] * len(batch),   # ← MUST match input length
        })

        for item in sorted(data["data"], key=lambda x: x["index"]):
            all_vecs.append(normalize(item["embedding"]))

        log.info(f"Batch {batch_num}/{n_batches} done ✓")

    return all_vecs


def embed_one(text: str, mode="query") -> list:
    log.info(f"Embedding query ({mode}): '{text[:70]}...'")
    return embed_texts([text], mode)[0]


def get_embed_dim() -> int:
    return len(embed_one("probe", "passage"))

# =============================================================================
# MILVUS SCHEMA
# =============================================================================

_FIELDS = {"id","vector","doc_id","file_name","file_type","modality",
           "page","section","parent_id","ingested_at","content"}

def _schema_valid(db):
    try:
        return _FIELDS.issubset({f["name"] for f in db.describe_collection(COLLECTION).get("fields",[])})
    except Exception:
        return False


def setup_collection(db, force=False):
    if db.has_collection(COLLECTION):
        if not force and _schema_valid(db):
            log.info("Collection exists with correct schema")
            return {"status": "already exists"}
        log.info("Dropping stale collection...")
        db.drop_collection(COLLECTION)

    log.info("Probing embed dimension...")
    dim    = get_embed_dim()
    schema = db.create_schema(auto_id=True, enable_dynamic_field=False)
    schema.add_field("id",          DataType.INT64,        is_primary=True)
    schema.add_field("vector",      DataType.FLOAT_VECTOR, dim=dim)
    schema.add_field("doc_id",      DataType.VARCHAR,      max_length=32)
    schema.add_field("file_name",   DataType.VARCHAR,      max_length=256)
    schema.add_field("file_type",   DataType.VARCHAR,      max_length=16)
    schema.add_field("modality",    DataType.VARCHAR,      max_length=16)
    schema.add_field("page",        DataType.INT64)
    schema.add_field("section",     DataType.VARCHAR,      max_length=256)
    schema.add_field("parent_id",   DataType.VARCHAR,      max_length=32)
    schema.add_field("ingested_at", DataType.VARCHAR,      max_length=32)
    schema.add_field("content",     DataType.VARCHAR,      max_length=5000)

    db.create_collection(COLLECTION, schema=schema)
    idx = db.prepare_index_params()
    idx.add_index("vector", index_type="FLAT", metric_type="COSINE")
    db.create_index(COLLECTION, idx)
    db.load_collection(COLLECTION)

    log.info(f"Collection '{COLLECTION}' created — dim={dim}")
    return {"status": "created", "dim": dim}

# =============================================================================
# RERANKER
# =============================================================================

def rerank(query: str, passages: list) -> tuple:
    if not passages:
        return passages, 0.0
    log.info(f"Reranking {len(passages)} passages...")
    try:
        data   = _post("ranking", {"model": RERANK_MODEL, "query": query,
                                   "passages": passages}, timeout=(5, 30))
        ranked = sorted(data["results"], key=lambda x: x["score"], reverse=True)
        result = [passages[r["index"]] for r in ranked]
        avg    = float(np.mean([r["score"] for r in ranked]))
        log.info(f"Rerank done — avg={avg:.4f}")
        return result, avg
    except Exception as e:
        log.warning(f"Reranker failed ({e}) — using vector order")
        return passages, 0.0

# =============================================================================
# DEDUP + TRIM
# =============================================================================

def deduplicate(passages: list, vectors: list) -> list:
    kept_t, kept_v = [], []
    for txt, vec in zip(passages, vectors):
        if not any(cosine(vec, kv) > DEDUP_THRESH for kv in kept_v):
            kept_t.append(txt)
            kept_v.append(vec)
    removed = len(passages) - len(kept_t)
    if removed:
        log.info(f"Dedup removed {removed} near-duplicate chunks")
    return kept_t


def trim_budget(passages: list) -> list:
    out, total = [], 0
    for p in passages:
        tc = token_count(p)
        if total + tc > TOKEN_BUDGET:
            break
        out.append(p)
        total += tc
    log.info(f"Token trim: {len(passages)} → {len(out)} chunks ({total} tokens)")
    return out or passages[:1]

# =============================================================================
# LLM
# =============================================================================

def call_llm(model: str, query: str, context: str, temp: float) -> str:
    log.info(f"LLM call: {model} temp={temp}")
    data = _post("chat/completions", {
        "model": model,
        "messages": [
            {"role": "system", "content": (
                "You are a precise document assistant. "
                "Answer ONLY from the provided context. "
                "Reproduce equations, formulas, and tables exactly as they appear in context. "
                "If the answer is not in context, say so. Never hallucinate."
            )},
            {"role": "user",
             "content": f"Context:\n{context}\n\nQuestion: {query}"}
        ],
        "temperature": temp,
        "max_tokens":  800,
    }, timeout=(5, 90))
    answer = data["choices"][0]["message"]["content"].strip()
    log.info(f"LLM response: {len(answer)} chars, model={model}")
    return answer


def llm_with_fallback(query: str, context: str, temp: float) -> tuple:
    for attempt in range(2):
        try:
            return call_llm(PRIMARY_LLM, query, context, temp), PRIMARY_LLM
        except Exception as e:
            log.warning(f"Primary LLM attempt {attempt+1} failed: {e}")
            if attempt == 0:
                time.sleep(1)
    log.warning("Falling back to secondary LLM...")
    try:
        return call_llm(FALLBACK_LLM, query, context, temp), FALLBACK_LLM
    except Exception as e:
        raise HTTPException(502, f"Both LLMs failed: {e}")

# =============================================================================
# CONFIDENCE SCORE
# =============================================================================

def confidence_score(vec_score: float, rerank_score: float, n_chunks: int) -> dict:
    """
    40% vector similarity  — how close retrieved chunks are to query
    40% reranker score     — semantic relevance confirmation
    20% context density    — how many useful chunks were found
    """
    density = min(1.0, n_chunks / 10.0)
    score   = round(0.4 * vec_score + 0.4 * rerank_score + 0.2 * density, 4)
    label   = "HIGH" if score > 0.7 else "MEDIUM" if score > 0.4 else "LOW"
    log.info(f"Confidence: {score} [{label}] vec={vec_score:.3f} rerank={rerank_score:.3f} density={density:.2f}")
    return {
        "score":         score,
        "label":         label,
        "vec_score":     round(vec_score, 4),
        "rerank_score":  round(rerank_score, 4),
    }

# =============================================================================
# APP
# =============================================================================

db_client: Optional[MilvusClient] = None

@asynccontextmanager
async def lifespan(app: FastAPI):
    global db_client

    log.info("=" * 60)
    log.info("Multimodal RAG API v3 — starting up")
    log.info(f"OCR       : {'PaddleOCR active' if OCR_AVAILABLE else 'not available'}")
    log.info(f"Embed model: {EMBED_MODEL}")
    log.info(f"LLM       : {PRIMARY_LLM} → fallback {FALLBACK_LLM}")

    # NVIDIA connectivity probe
    try:
        r = session.post(
            f"{BASE_URL}/embeddings",
            headers={"Authorization": f"Bearer {NVIDIA_KEY}"},
            json={"model": EMBED_MODEL, "input": ["ping"],
                  "input_type": "passage", "encoding_format": "float",
                  "modality": ["text"]},   # single input → single modality
            timeout=(4, 15),
        )
        log.info(f"NVIDIA API reachable ✓  (status={r.status_code})")
    except Exception as e:
        log.warning(f"NVIDIA API unreachable at startup: {e}")

    db_client = MilvusClient(DB_PATH)
    log.info(f"Milvus Lite connected → {DB_PATH}")
    log.info("=" * 60)

    yield

    log.info("Shutdown complete")


app = FastAPI(
    title="Multimodal RAG API",
    version="3.0",
    description="PDF · DOCX · PPTX · HTML ingestion with OCR, reranking, and LLM answering",
    lifespan=lifespan,
)

# =============================================================================
# /create
# =============================================================================

@app.post("/create")
def create(force: bool = False):
    """
    Create the Milvus vector collection.
    Pass **?force=true** to drop and recreate (required after schema changes).
    """
    log.info(f"/create — force={force}")
    result = setup_collection(db_client, force=force)
    return result

# =============================================================================
# /ingest
# =============================================================================

@app.post("/ingest")
async def ingest(file: UploadFile = File(...)):
    """
    Ingest a document.

    - **Text** → sentence-chunked with 15% overlap
    - **Tables** → converted to structured text rows
    - **Images** → OCR via PaddleOCR → text chunk (skips tiny < 5KB images)

    All chunks batch-embedded in ~1–3 API calls (not one per chunk).
    """
    log.info(f"/ingest — file='{file.filename}'")

    if not db_client.has_collection(COLLECTION):
        raise HTTPException(400, "Run /create first.")

    raw    = await file.read()
    doc_id = hashlib.md5(raw).hexdigest()[:16]
    log.info(f"File size: {len(raw):,} bytes, doc_id={doc_id}")

    chunks = parse_file(raw, file.filename, doc_id)
    if not chunks:
        raise HTTPException(422, "No content extracted from file.")

    n_text  = sum(1 for c in chunks if c.modality == "text")
    n_table = sum(1 for c in chunks if c.modality == "table")
    n_img   = sum(1 for c in chunks if c.modality == "image_ocr")
    log.info(f"Chunks: {len(chunks)} total — text={n_text} table={n_table} image_ocr={n_img}")

    # Batch embed all
    texts = [c.text[:MAX_CONTENT] for c in chunks]
    log.info(f"Starting batch embedding: {len(texts)} chunks, batch_size={EMBED_BATCH}")
    vecs  = embed_texts(texts, mode="passage")

    # Build + insert records
    ts    = now()
    records = [{
        "vector":      vec,
        "doc_id":      c.doc_id,
        "file_name":   c.file_name,
        "file_type":   c.file_type,
        "modality":    c.modality,
        "page":        c.page,
        "section":     (c.section or "")[:255],
        "parent_id":   c.parent_id or "",
        "ingested_at": ts,
        "content":     c.text[:MAX_CONTENT],
    } for c, vec in zip(chunks, vecs)]

    inserted = 0
    for i in range(0, len(records), 100):
        db_client.insert(COLLECTION, records[i:i+100])
        inserted += len(records[i:i+100])
        log.info(f"Inserted {inserted}/{len(records)}")

    log.info(f"/ingest complete — {len(records)} records stored for '{file.filename}'")
    return {
        "status":    "ok",
        "doc_id":    doc_id,
        "file":      file.filename,
        "file_type": chunks[0].file_type,
        "chunks":    len(records),
        "breakdown": {"text": n_text, "table": n_table, "image_ocr": n_img},
        "ocr":       OCR_AVAILABLE,
    }

# =============================================================================
# /ask
# =============================================================================

@app.post("/ask")
def ask(
    query:       str   = Query(...),
    temperature: float = Query(0.2, ge=0.0, le=1.0,
                               description="0.1=exact/equations · 0.2=factual · 0.3=balanced · 0.5=creative"),
    top_k:       int   = Query(20,  ge=1,   le=100),
    file_name:   Optional[str] = Query(None,
                               description="Filter to a specific ingested file"),
):
    """
    Query the RAG pipeline.

    **Temperature guide:**
    - 0.1 → retrieve equations, formulas, exact data
    - 0.2 → factual Q&A (default)
    - 0.3 → balanced reasoning
    - 0.5 → creative synthesis / summaries

    **Pipeline:** embed → vector search → rerank → dedup → token trim → LLM → confidence
    """
    log.info(f"/ask — '{query[:80]}' temp={temperature} top_k={top_k} filter={file_name}")

    if not query.strip():
        raise HTTPException(400, "Query is empty.")
    if not db_client.has_collection(COLLECTION):
        raise HTTPException(400, "Run /create and /ingest first.")

    # 1. Embed query
    q_vec = embed_one(query, "query")

    # 2. Vector search
    expr = f'file_name == "{file_name}"' if file_name else None
    results = db_client.search(
        COLLECTION, data=[q_vec], limit=top_k, filter=expr,
        search_params={"metric_type": "COSINE"},
        output_fields=["content","vector","modality","file_name","page","section","doc_id"],
    )

    if not results or not results[0]:
        raise HTTPException(404, "No results found — ingest documents first.")

    hits          = results[0]
    passages      = [h["entity"]["content"] for h in hits]
    hit_vectors   = [h["entity"].get("vector", q_vec) for h in hits]
    top_vec_score = float(hits[0].get("distance", 0.0))

    log.info(f"Vector search: {len(hits)} hits, top_score={top_vec_score:.4f}")
    log.info(f"Modalities: {set(h['entity'].get('modality') for h in hits)}")

    sources = [{
        "file":     h["entity"].get("file_name",""),
        "page":     h["entity"].get("page", 0),
        "section":  h["entity"].get("section",""),
        "modality": h["entity"].get("modality",""),
        "score":    round(float(h.get("distance", 0.0)), 4),
    } for h in hits[:5]]

    # 3. Rerank
    reranked, rerank_score = rerank(query, passages)
    reranked    = reranked[:10]
    rerank_vecs = hit_vectors[:min(10, len(hit_vectors))]

    # 4. Dedup
    deduped = deduplicate(reranked, rerank_vecs)

    # 5. Token trim
    final = trim_budget(deduped)

    # 6. LLM
    context = "\n\n---\n\n".join(final)
    log.info(f"LLM input: {len(final)} chunks, {token_count(context)} tokens")
    answer, model_used = llm_with_fallback(query, context, temperature)

    # 7. Confidence
    conf = confidence_score(top_vec_score, rerank_score, len(final))

    log.info(f"/ask done — model={model_used} confidence={conf['score']} [{conf['label']}]")

    return {
        "answer":      answer,
        "confidence":  conf,
        "model":       model_used,
        "temperature": temperature,
        "chunks_used": len(final),
        "top_sources": sources[:3],
    }